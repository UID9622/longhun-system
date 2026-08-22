#!/usr/bin/env python3
# DNA: #龍芯⚡️丙午·丙申·庚戌·壬午·䷙大畜-SCRIPT-MANAGER-v1.2-UID9622
# CONFIRM: #CONFIRM🌌9622-ONLY-ONCE🧬LK9X-772Z
# SEAL: #ZHUGEXIN⚡️2025-🇨🇳🐉⚖️♠️🧚🏼‍♀️❤️♾️-DEVICE-BIND-SOUL
#龍芯⚡️丙午·丙申·癸酉·庚申·䷒临-MAKE_PPT-v1.0-ed8bcff8
# CREATOR: 诸葛鑫 (UID9622)
# PROTOCOL: CC BY-NC-SA 4.0
# -*- coding: utf-8 -*-
"""
龍魂·PPT生成器 (python-pptx)
依赖: 需激活 venv (.venv_docs) 后 python-pptx 可用
功能: 用龍魂暗色金模板，从结构化数据生成 .pptx。

用法:
    source .venv_docs/bin/activate
    python3 bin/docs_factory/make_ppt.py --title "标题" --out output/ppt/demo.pptx
"""
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor

DRAGON_RED = RgbColor(196, 30, 58)
DRAGON_GOLD = RgbColor(212, 165, 116)
DRAGON_BLACK = RgbColor(26, 26, 26)
GREY = RgbColor(128, 128, 128)
DNA = "#CONFIRM🌌9622-ONLY-ONCE🧬LK9X-772Z"

def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_bg(slide, prs):
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    _solid(bg, DRAGON_BLACK)
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg

def build(title, subtitle, out):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 标题页
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    tb = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54); p.font.bold = True; p.font.color.rgb = DRAGON_RED
    sub = s.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1))
    p = sub.text_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28); p.font.color.rgb = DRAGON_GOLD
    dna = s.shapes.add_textbox(Inches(1), Inches(6.7), Inches(11), Inches(0.5))
    p = dna.text_frame.paragraphs[0]
    p.text = "DNA: " + DNA
    p.font.size = Pt(10); p.font.color.rgb = GREY

    prs.save(out)
    return f"PPT已生成: {out}"

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--title", default="🐉 龍魂系统")
    ap.add_argument("--subtitle", default="技术为人民 · 不为资本")
    ap.add_argument("--out", default="output/ppt/longhun.pptx")
    a = ap.parse_args()
    print(build(a.title, a.subtitle, a.out))

if __name__ == "__main__":
    main()
